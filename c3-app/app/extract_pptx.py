"""
Extract text content from PowerPoint presentation
"""
from pptx import Presentation
import json
import sys

def extract_pptx_content(pptx_path):
    """Extract text content from a PowerPoint file"""
    try:
        prs = Presentation(pptx_path)
        slides_content = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_data = {
                'slide_number': slide_num,
                'title': '',
                'content': []
            }
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # First text box is usually the title
                    if not slide_data['title'] and shape.shape_type == 1:  # Text box
                        slide_data['title'] = text
                    else:
                        slide_data['content'].append(text)
            
            slides_content.append(slide_data)
        
        return slides_content
    except Exception as e:
        print(f"Error reading PowerPoint: {e}", file=sys.stderr)
        return None

if __name__ == "__main__":
    pptx_path = "/Users/ck/Downloads/C3Ops_Pricing_Presentation.pptx"
    
    content = extract_pptx_content(pptx_path)
    
    if content:
        print(json.dumps(content, indent=2))
    else:
        sys.exit(1)
